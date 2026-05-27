"""PowerPoint 解析器"""
from pptx import Presentation


def parse_pptx(file_path: str) -> list[dict]:
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            pages.append({"page": i + 1, "text": "\n".join(texts)})
    return pages
